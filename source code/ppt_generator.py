##Example PPT Code

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import datetime

def create_ppt(company_name, data, start_date_str, end_date_str):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = f"{company_name} 데이터 보고서"
    subtitle.text = f"기간: {start_date_str} ~ {end_date_str}"

    # 데이터 슬라이드 추가
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "데이터 그래프"

    # 차트 데이터 설정
    chart_data = CategoryChartData()
    chart_data.categories = [d[-1].date() if isinstance(d[-1], datetime.datetime) else datetime.datetime.strptime(d[-1], '%Y-%m-%d').date() for d in data]
    chart_data.add_series('Series 1', (d[1] for d in data))

    # 차트 추가
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    # 파일 저장
    pptx_file = f"{company_name}_report_{start_date_str}_{end_date_str}.pptx"
    prs.save(pptx_file)
    return pptx_file
    return pptx_file